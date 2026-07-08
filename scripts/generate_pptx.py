#!/usr/bin/env python3
"""Generate editable PPTX from the 投喂版 markdown using python-pptx.

Usage:
    python3 generate_pptx.py <投喂版.md路径> <输出.pptx路径> [页数]

Output:
    A 16:9 .pptx file with editable text boxes for each slide.
    Cover/closing slides use IKB blue background; content slides alternate white/dark.
    All text is in independent text boxes, editable in PowerPoint/WPS/Keynote.
"""

import os, sys, re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Colors
IKB_BLUE   = RGBColor(0x00, 0x2F, 0xA7)
DARK_TEXT   = RGBColor(0x0A, 0x0A, 0x0A)
GREY_TEXT   = RGBColor(0x52, 0x52, 0x52)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG     = RGBColor(0x0A, 0x0A, 0x0A)
DARK_BODY   = RGBColor(0xCC, 0xCC, 0xCC)
LIGHT_META  = RGBColor(0x99, 0xAA, 0xDD)


def parse_feed_markdown(path: str) -> list[dict]:
    """Parse the 投喂版 markdown into structured slide data."""
    with open(path, "r") as f:
        content = f.read()

    pages_raw = re.split(r'\n---\n', content)
    slides = []
    in_notes = False

    for part in pages_raw:
        if "# 讲标人备注" in part:
            in_notes = True
            continue
        if in_notes:
            continue
        if part.strip().startswith("## 第") and "页 ·" in part:
            lines = part.strip().split("\n")
            title = ""
            bullets = []
            kpi = ""

            for line in lines:
                line = line.strip()
                if not line or line.startswith("**项目**") or line.startswith("**讲标"):
                    continue
                if line.startswith("## 第") and "页 ·" in line:
                    title = re.sub(r"（.*?布局）", "", line.split("·")[-1].strip()).strip()
                    continue
                if line.startswith("**幻灯片展示内容**"):
                    continue
                if line.startswith("- **数据大字报**：") or line.startswith("- **数据大字报**: "):
                    kpi = line.replace("- **数据大字报**：", "").replace("- **数据大字报**: ", "").replace("**", "").strip()
                    continue
                if line.startswith("- ") or line.startswith("1. ") or line.startswith("2. ") or line.startswith("3. ") or line.startswith("4. ") or line.startswith("5. ") or line.startswith("6. ") or line.startswith("7. "):
                    bullet = re.sub(r'^[\d]+\.\s*\*\*|^-\s*\*\*|^[\d]+\.\s*|^-\s*', '', line).strip().replace("**", "")
                    if bullet and len(bullet) > 5:
                        bullets.append(bullet)
                    continue
                if bullets and line and not line.startswith("#"):
                    bullets[-1] += " " + line

            slides.append({"title": title, "bullets": bullets, "kpi": kpi})

    return slides


def add_textbox(slide, left, top, width, height, text, font_size=Pt(18),
                color=DARK_TEXT, bold=False, font_name="Microsoft YaHei",
                alignment=None, line_spacing=1.3):
    """Add a single-line text box to a slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    if alignment is not None:
        p.alignment = alignment
    p.space_after = Pt(4)
    p.line_spacing = line_spacing
    return tf


def generate_pptx(feed_path: str, output_path: str):
    """Generate PPTX from feed markdown."""
    slides_data = parse_feed_markdown(feed_path)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    dark_pages = {4, 7, 9, 12, 14, 16}  # Page numbers with dark theme

    for i, slide_data in enumerate(slides_data):
        page_num = i + 1
        title = slide_data["title"]
        bullets = slide_data["bullets"]
        kpi = slide_data["kpi"]

        is_cover = (page_num == 1)
        is_closing = (page_num == len(slides_data))
        is_dark = page_num in dark_pages

        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)

        if is_cover:
            bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = IKB_BLUE
            add_textbox(slide, 1.5, 1.8, 10.3, 2.0, "广汇能源安全生产信息化\n管理平台（二期）",
                        Pt(48), WHITE, line_spacing=1.05)
            add_textbox(slide, 1.5, 4.2, 10.3, 0.6, "技术响应方案讲标", Pt(22), RGBColor(0xCC, 0xCC, 0xDD))
            add_textbox(slide, 1.5, 5.2, 10.3, 0.6, kpi, Pt(18), LIGHT_META)
            add_textbox(slide, 1.5, 6.3, 10.3, 0.4,
                        "讲标团队：商务负责人 + 核心技术负责人  ·  2026 年 7 月  ·  哈密淖毛湖镇",
                        Pt(12), RGBColor(0x88, 0x99, 0xBB))

        elif is_closing:
            bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = IKB_BLUE
            add_textbox(slide, 1.5, 2.0, 10.3, 1.2, "感谢各位评委聆听", Pt(44), WHITE)
            add_textbox(slide, 1.5, 3.5, 10.3, 0.6, "广汇能源安全生产信息化管理平台（二期）技术方案讲标",
                        Pt(20), RGBColor(0xCC, 0xCC, 0xDD))
            add_textbox(slide, 1.5, 4.5, 10.3, 0.6,
                        "我们准备好了。期待与广汇能源携手，共同打造煤化工安全生产数字化的标杆工程。",
                        Pt(18), LIGHT_META)
            if kpi:
                add_textbox(slide, 1.5, 5.8, 10.3, 0.6, kpi, Pt(18), LIGHT_META)

        else:
            bg_color = DARK_BG if is_dark else WHITE
            text_color = WHITE if is_dark else DARK_TEXT
            body_color = DARK_BODY if is_dark else GREY_TEXT
            bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = bg_color

            # Accent line
            line = slide.shapes.add_shape(1, Inches(0.8), Inches(0.55), Inches(0.06), Inches(0.6))
            line.fill.solid(); line.fill.fore_color.rgb = IKB_BLUE; line.line.fill.background()

            # Title
            add_textbox(slide, 1.1, 0.5, 10.5, 0.8, title, Pt(28), text_color, line_spacing=1.1)

            # KPI (top right)
            if kpi:
                kpi_color = IKB_BLUE if not is_dark else RGBColor(0x88, 0xAA, 0xFF)
                add_textbox(slide, 8.5, 0.4, 4.2, 0.9, kpi, Pt(14), kpi_color, bold=True,
                            alignment=2, line_spacing=1.2)  # alignment=2 = right

            # Bullets
            if bullets:
                font_sz = Pt(13) if len(bullets) > 4 else Pt(14) if len(bullets) > 3 else Pt(15)
                txBox = slide.shapes.add_textbox(Inches(1.1), Inches(1.6), Inches(11.5), Inches(5.3))
                tf = txBox.text_frame; tf.word_wrap = True
                for j, bullet in enumerate(bullets):
                    p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                    p.text = "• " + bullet[:400] + ("…" if len(bullet) > 400 else "")
                    p.font.size = font_sz; p.font.color.rgb = body_color
                    p.font.name = "Microsoft YaHei"; p.line_spacing = 1.35 if j == 0 else 1.25
                    p.space_after = Pt(6)

            # Page number
            pn_color = RGBColor(0x99, 0x99, 0x99) if not is_dark else RGBColor(0x66, 0x66, 0x66)
            add_textbox(slide, 11.8, 7.0, 1.2, 0.3, f"{page_num} / {len(slides_data)}",
                        Pt(10), pn_color, alignment=2)

    prs.save(output_path)
    return len(prs.slides)


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python3 generate_pptx.py <投喂版.md> <输出.pptx>")
        sys.exit(1)
    n = generate_pptx(sys.argv[1], sys.argv[2])
    print(f"Done: {n} slides → {sys.argv[2]}")
